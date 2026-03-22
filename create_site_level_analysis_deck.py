"""
Build a PowerPoint deck for the Site-Level Analysis dashboard with chart PNGs embedded.

Tries Plotly Kaleido first; if unavailable or broken, redraws the same data with Matplotlib
so charts still appear (typical on locked-down corporate laptops).

  pip install python-pptx plotly numpy matplotlib
  pip install kaleido   # optional, higher-fidelity match to Streamlit

Output: Site_Level_Analysis_Dashboard.pptx
"""
from __future__ import annotations

import tempfile
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from deck_plotly_export import write_figure_png, write_sparkline_png, write_wide_spark_png
from site_level_charts import (
    GRAY,
    GREEN,
    ORANGE,
    PINK,
    PINK_DEEP,
    fig_region_avail_bars,
    fig_sparkline,
    fig_sparkline_green,
    fig_sparkline_orange,
    fig_sparkline_pink,
    fig_stacked_h_bar,
)

MAGENTA = RGBColor(233, 69, 96)
DARK_BG = RGBColor(26, 26, 46)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 200)
CARD_BG = RGBColor(40, 40, 60)
RGB_GREEN = RGBColor(0, 217, 165)
RGB_ORANGE = RGBColor(247, 127, 0)
RGB_MAGENTA = RGBColor(233, 69, 96)


def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    return slide


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = _blank_slide(prs)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.333), Inches(1.4))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = MAGENTA
    p.alignment = PP_ALIGN.CENTER

    sb = slide.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(12.333), Inches(1))
    p = sb.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(22)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    db = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    p = db.text_frame.paragraphs[0]
    p.text = datetime.now().strftime("%B %d, %Y")
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER


def add_section_slide(prs: Presentation, tab_name: str) -> None:
    slide = _blank_slide(prs)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(12.333), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = tab_name
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = MAGENTA
    p.alignment = PP_ALIGN.CENTER
    sb = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.6))
    p2 = sb.text_frame.paragraphs[0]
    p2.text = "Same Plotly figure definitions as the Streamlit app (PNG via Kaleido or Matplotlib)"
    p2.font.size = Pt(16)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], footer: str | None = None) -> None:
    slide = _blank_slide(prs)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.75))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = MAGENTA

    bb = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(11.8), Inches(5.6))
    tf = bb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = f"• {line}"
        para.font.size = Pt(20)
        para.font.color.rgb = WHITE
        para.space_after = Pt(14)

    if footer:
        fb = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(12.333), Inches(0.45))
        p = fb.text_frame.paragraphs[0]
        p.text = footer
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY


def add_metrics_slide(prs: Presentation, title: str, metrics: list[tuple[str, str, RGBColor]]) -> None:
    slide = _blank_slide(prs)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.75))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = MAGENTA

    card_w, card_h = Inches(3.85), Inches(1.85)
    gap = Inches(0.28)
    sx, sy = Inches(0.5), Inches(1.25)

    for i, (label, value, color) in enumerate(metrics):
        row, col = divmod(i, 3)
        x = sx + col * (card_w + gap)
        y = sy + row * (card_h + gap)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(2)

        vb = slide.shapes.add_textbox(x, y + Inches(0.25), card_w, Inches(0.95))
        vp = vb.text_frame.paragraphs[0]
        vp.text = str(value)
        vp.font.size = Pt(36)
        vp.font.bold = True
        vp.font.color.rgb = color
        vp.alignment = PP_ALIGN.CENTER

        lb = slide.shapes.add_textbox(x, y + Inches(1.15), card_w, Inches(0.55))
        lp = lb.text_frame.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(14)
        lp.font.color.rgb = LIGHT_GRAY
        lp.alignment = PP_ALIGN.CENTER


def add_chart_slide(prs: Presentation, tab: str, caption: str, fig, tmp: Path, seq: list[int]) -> None:
    seq[0] += 1
    png = tmp / f"chart_{seq[0]}.png"
    ly = fig.layout.height
    base_h = int(ly) + 90 if ly is not None else 400
    h = min(max(base_h, 260), 520)
    ok = write_figure_png(fig, png, width=1320, height=h)

    slide = _blank_slide(prs)
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.28), Inches(12.4), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{tab} — {caption}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = MAGENTA

    if ok:
        slide.shapes.add_picture(str(png), Inches(0.35), Inches(0.95), width=Inches(12.6))
    else:
        fb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.2))
        fp = fb.text_frame.paragraphs[0]
        fp.text = "Chart image export failed (install kaleido and/or matplotlib)."
        fp.font.size = Pt(18)
        fp.font.color.rgb = LIGHT_GRAY


def add_sparkline_grid_slide(prs: Presentation, tab: str, tmp: Path, seq: list[int]) -> None:
    labels = [
        "KPI 1 (green trend)",
        "KPI 2 (green trend)",
        "KPI 3 (orange trend)",
        "KPI 4 (orange trend)",
        "KPI 5 (pink trend)",
        "KPI 6 (pink trend)",
    ]
    fns = [
        fig_sparkline_green,
        fig_sparkline_green,
        fig_sparkline_orange,
        fig_sparkline_orange,
        fig_sparkline_pink,
        fig_sparkline_pink,
    ]

    slide = _blank_slide(prs)
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.28), Inches(12.4), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{tab} — KPI sparklines (six-pack, matches dashboard row)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = MAGENTA

    col_w = Inches(3.95)
    row_h_img = Inches(1.35)
    gap_x = Inches(0.22)
    gap_y = Inches(0.55)
    x0, y0 = Inches(0.4), Inches(0.95)

    for i, (fn, lbl) in enumerate(zip(fns, labels)):
        seq[0] += 1
        png = tmp / f"spark_{seq[0]}.png"
        f = fn()
        f.update_layout(height=90, margin=dict(l=8, r=8, t=4, b=4))
        ok = write_sparkline_png(f, png, width=520, height=200)

        row, col = divmod(i, 3)
        left = x0 + col * (col_w + gap_x)
        top = y0 + row * (row_h_img + gap_y + Inches(0.38))

        if ok:
            slide.shapes.add_picture(str(png), left, top, width=col_w)

        lb = slide.shapes.add_textbox(left, top + row_h_img, col_w, Inches(0.32))
        lp = lb.text_frame.paragraphs[0]
        lp.text = lbl
        lp.font.size = Pt(11)
        lp.font.color.rgb = LIGHT_GRAY


def add_wide_spark_slide(prs: Presentation, tab: str, caption: str, fig, tmp: Path, seq: list[int]) -> None:
    seq[0] += 1
    png = tmp / f"wide_{seq[0]}.png"
    fig.update_layout(height=220, margin=dict(l=20, r=20, t=40, b=30))
    ok = write_wide_spark_png(fig, png)

    slide = _blank_slide(prs)
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.28), Inches(12.4), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{tab} — {caption}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = MAGENTA
    if ok:
        slide.shapes.add_picture(str(png), Inches(0.45), Inches(1.0), width=Inches(12.4))
    else:
        fb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(0.8))
        fp = fb.text_frame.paragraphs[0]
        fp.text = "Trend chart export failed."
        fp.font.size = Pt(18)
        fp.font.color.rgb = LIGHT_GRAY


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    seq = [0]

    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)

        add_title_slide(
            prs,
            "Site-Level Analysis Dashboard",
            "Streamlit views + Plotly chart exports (all tabs)",
        )

        add_bullet_slide(
            prs,
            "About this deck",
            [
                "Following slides mirror each Streamlit tab: KPI sparkline grid, then every bar chart from that tab.",
                "Figures are generated with the shared module site_level_charts.py (same logic as the live app).",
                "PNG export: Kaleido when available, otherwise Matplotlib redraw (pip install matplotlib).",
            ],
        )

        # --- Executive Summary ---
        add_section_slide(prs, "Executive Summary")
        add_metrics_slide(
            prs,
            "Executive Summary — reference KPIs",
            [
                ("Network availability", "99.58%", RGB_GREEN),
                ("Downtime (enterprise)", "412.7M sec", RGB_GREEN),
                ("Outage events", "2,847", RGB_ORANGE),
                ("Outage minutes", "428K", RGB_ORANGE),
                ("Customer minutes", "18.2M", RGB_MAGENTA),
                ("Impacted subscribers", "512K", RGB_MAGENTA),
            ],
        )
        add_sparkline_grid_slide(prs, "Executive Summary", tmp, seq)
        add_chart_slide(
            prs,
            "Executive Summary",
            "Availability — Summary categories",
            fig_stacked_h_bar(
                "Availability — Summary categories",
                [("Transport", 61, GRAY), ("RAN", 26, PINK), ("Power", 12, PINK_DEEP)],
                height=140,
            ),
            tmp,
            seq,
        )
        add_chart_slide(
            prs,
            "Executive Summary",
            "COTTR — Summary categories",
            fig_stacked_h_bar(
                "COTTR — Summary categories",
                [("Transport", 76, GRAY), ("RAN", 16, PINK), ("Other", 8, "#4b5563")],
                height=140,
            ),
            tmp,
            seq,
        )

        # --- Site Analysis ---
        add_section_slide(prs, "Site Analysis")
        add_sparkline_grid_slide(prs, "Site Analysis", tmp, seq)
        add_chart_slide(
            prs,
            "Site Analysis",
            "Availability — Summary categories",
            fig_stacked_h_bar(
                "Availability — Summary categories",
                [("Transport", 61, GRAY), ("RAN", 26, PINK), ("Power", 12, PINK_DEEP)],
            ),
            tmp,
            seq,
        )
        add_chart_slide(
            prs,
            "Site Analysis",
            "Availability — Focus categories",
            fig_stacked_h_bar(
                "Availability — Focus categories",
                [
                    ("Transport-AAV", 54, GRAY),
                    ("Site Mod", 22, PINK),
                    ("Maintenance", 8, PINK_DEEP),
                    ("Decommission", 7, "#7f1d1d"),
                    ("Other", 9, "#4b5563"),
                ],
            ),
            tmp,
            seq,
        )
        add_chart_slide(
            prs,
            "Site Analysis",
            "COTTR — Summary categories",
            fig_stacked_h_bar(
                "COTTR — Summary categories",
                [("Transport", 76, GRAY), ("RAN", 16, PINK), ("Other", 8, "#4b5563")],
            ),
            tmp,
            seq,
        )
        add_chart_slide(
            prs,
            "Site Analysis",
            "COTTR — Focus categories",
            fig_stacked_h_bar(
                "COTTR — Focus categories",
                [
                    ("Transport-AAV", 64, GRAY),
                    ("Site Mod", 12, PINK),
                    ("Maintenance", 10, PINK_DEEP),
                    ("Decommission", 6, "#7f1d1d"),
                    ("Other", 8, "#4b5563"),
                ],
            ),
            tmp,
            seq,
        )

        # --- Region Availability ---
        add_section_slide(prs, "Region Availability")
        add_sparkline_grid_slide(prs, "Region Availability", tmp, seq)
        add_chart_slide(
            prs,
            "Region Availability",
            "Availability % by region (weighted)",
            fig_region_avail_bars(
                "Trailing 30 days (weighted)",
                [
                    ("South", 99.612),
                    ("Mid-Atlantic", 99.523),
                    ("Northeast", 99.541),
                    ("West", 99.498),
                    ("Central", 99.455),
                ],
                height=300,
            ),
            tmp,
            seq,
        )
        add_chart_slide(
            prs,
            "Region Availability",
            "South — summary categories",
            fig_stacked_h_bar(
                "South — summary",
                [("Transport", 58, GRAY), ("RAN", 28, PINK), ("Power", 14, PINK_DEEP)],
            ),
            tmp,
            seq,
        )
        add_chart_slide(
            prs,
            "Region Availability",
            "Northeast — summary categories",
            fig_stacked_h_bar(
                "Northeast — summary",
                [("Transport", 64, GRAY), ("RAN", 24, PINK), ("Power", 12, PINK_DEEP)],
            ),
            tmp,
            seq,
        )

        # --- Area Availability ---
        add_section_slide(prs, "Area Availability")
        add_sparkline_grid_slide(prs, "Area Availability", tmp, seq)
        add_chart_slide(
            prs,
            "Area Availability",
            "Availability % by market (West)",
            fig_region_avail_bars(
                "Top markets by population weight",
                [
                    ("Houston", 99.582),
                    ("Dallas", 99.541),
                    ("Phoenix", 99.498),
                    ("Denver", 99.455),
                    ("Seattle", 99.521),
                    ("Portland", 99.389),
                ],
                height=320,
            ),
            tmp,
            seq,
        )
        add_chart_slide(
            prs,
            "Area Availability",
            "Availability — focus",
            fig_stacked_h_bar(
                "Availability — focus",
                [
                    ("Transport-AAV", 52, GRAY),
                    ("Site Mod", 24, PINK),
                    ("Maintenance", 9, PINK_DEEP),
                    ("Decommission", 6, "#7f1d1d"),
                    ("Other", 9, "#4b5563"),
                ],
            ),
            tmp,
            seq,
        )
        add_chart_slide(
            prs,
            "Area Availability",
            "COTTR — focus",
            fig_stacked_h_bar(
                "COTTR — focus",
                [
                    ("Transport-AAV", 62, GRAY),
                    ("Site Mod", 14, PINK),
                    ("Maintenance", 12, PINK_DEEP),
                    ("Decommission", 5, "#7f1d1d"),
                    ("Other", 7, "#4b5563"),
                ],
            ),
            tmp,
            seq,
        )

        # --- Availability Detail ---
        add_section_slide(prs, "Availability Detail")
        add_sparkline_grid_slide(prs, "Availability Detail", tmp, seq)
        for cap, segments in [
            (
                "Availability — summary",
                [("Transport", 61, GRAY), ("RAN", 26, PINK), ("Power", 12, PINK_DEEP)],
            ),
            (
                "Availability — focus",
                [
                    ("Transport-AAV", 54, GRAY),
                    ("Site Mod", 22, PINK),
                    ("Maintenance", 8, PINK_DEEP),
                    ("Decommission", 7, "#7f1d1d"),
                    ("Other", 9, "#4b5563"),
                ],
            ),
            (
                "Top drag — vendor slice",
                [("Vendor A", 44, GRAY), ("Vendor B", 31, PINK), ("Vendor C", 18, PINK_DEEP), ("Other", 7, "#4b5563")],
            ),
            (
                "Technology slice",
                [("LTE", 48, GRAY), ("5G NR", 35, PINK), ("Legacy", 17, PINK_DEEP)],
            ),
        ]:
            add_chart_slide(
                prs,
                "Availability Detail",
                cap,
                fig_stacked_h_bar(cap, segments),
                tmp,
                seq,
            )

        # --- COTTR Detail ---
        add_section_slide(prs, "COTTR Detail")
        add_sparkline_grid_slide(prs, "COTTR Detail", tmp, seq)
        for cap, segments in [
            ("COTTR — summary", [("Transport", 76, GRAY), ("RAN", 16, PINK), ("Other", 8, "#4b5563")]),
            (
                "COTTR — focus",
                [
                    ("Transport-AAV", 64, GRAY),
                    ("Site Mod", 12, PINK),
                    ("Maintenance", 10, PINK_DEEP),
                    ("Decommission", 6, "#7f1d1d"),
                    ("Other", 8, "#4b5563"),
                ],
            ),
            (
                "COTTR — by resolution",
                [("Remote", 38, GRAY), ("Dispatch", 34, PINK), ("Vendor", 21, PINK_DEEP), ("Other", 7, "#4b5563")],
            ),
            ("COTTR — priority", [("P1", 22, PINK_DEEP), ("P2", 41, PINK), ("P3", 28, GRAY), ("P4", 9, "#4b5563")]),
        ]:
            add_chart_slide(prs, "COTTR Detail", cap, fig_stacked_h_bar(cap, segments), tmp, seq)

        # --- Customer Minutes Detail ---
        add_section_slide(prs, "Customer Minutes Detail")
        add_sparkline_grid_slide(prs, "Customer Minutes Detail", tmp, seq)
        add_chart_slide(
            prs,
            "Customer Minutes Detail",
            "CM — summary drivers",
            fig_stacked_h_bar(
                "CM — summary drivers",
                [("Transport", 58, GRAY), ("RAN", 28, PINK), ("Power / other", 14, PINK_DEEP)],
            ),
            tmp,
            seq,
        )
        add_chart_slide(
            prs,
            "Customer Minutes Detail",
            "CM — focus drivers",
            fig_stacked_h_bar(
                "CM — focus drivers",
                [
                    ("Transport-AAV", 51, GRAY),
                    ("Site Mod", 21, PINK),
                    ("Maintenance", 11, PINK_DEEP),
                    ("Decommission", 9, "#7f1d1d"),
                    ("Other", 8, "#4b5563"),
                ],
            ),
            tmp,
            seq,
        )

        # --- Data Diagnostics ---
        add_section_slide(prs, "Data Diagnostics")
        add_sparkline_grid_slide(prs, "Data Diagnostics", tmp, seq)
        add_wide_spark_slide(
            prs,
            "Data Diagnostics",
            "Ingest volume trend (sample)",
            fig_sparkline(GREEN, n=36, seed=42),
            tmp,
            seq,
        )
        add_chart_slide(
            prs,
            "Data Diagnostics",
            "Rows by region (%)",
            fig_stacked_h_bar(
                "Rows by region (%)",
                [
                    ("South", 28, GRAY),
                    ("NE", 22, PINK),
                    ("West", 20, PINK_DEEP),
                    ("Central", 18, "#7f1d1d"),
                    ("Other", 12, "#4b5563"),
                ],
            ),
            tmp,
            seq,
        )
        add_chart_slide(
            prs,
            "Data Diagnostics",
            "Late arrivals (> 4h)",
            fig_stacked_h_bar(
                "Late arrivals (> 4h)",
                [("On time", 94, GREEN), ("Late", 6, PINK)],
            ),
            tmp,
            seq,
        )

        add_bullet_slide(
            prs,
            "Implementation",
            [
                "Streamlit app: app_site_level_analysis.py",
                "Shared Plotly builders: site_level_charts.py",
                "Deploy: Streamlit Community Cloud + GitHub streamlit-dashboards",
            ],
        )

        add_title_slide(prs, "Thank you", "Site-Level Analysis — charts & documentation")

        base = Path(__file__).resolve().parent / "Site_Level_Analysis_Dashboard.pptx"
        out = base
        try:
            prs.save(str(out))
        except OSError:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            out = base.with_name(f"Site_Level_Analysis_Dashboard_{ts}.pptx")
            prs.save(str(out))
            print(f"(Default file locked; saved alternate path.)")
        print(f"Saved: {out}")


if __name__ == "__main__":
    main()
